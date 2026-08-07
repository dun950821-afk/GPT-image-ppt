#!/usr/bin/env python3
"""Fail-closed validation for v2.1 refined editable PowerPoint deliveries."""

from __future__ import annotations

import argparse
import csv
import io
import json
import sys
import zipfile
from collections import defaultdict
from pathlib import Path
from typing import Dict, Iterable, List, Sequence, Tuple


MANIFEST_REQUIRED = {
    "slide_id",
    "element_id",
    "pptx_shape_name",
    "file_name",
    "relative_path",
    "element_type",
    "final_representation",
    "included_in_final",
    "decision_reason",
    "source_owner_id",
    "x",
    "y",
    "width",
    "height",
    "z_order",
    "opacity",
    "rotation",
    "source_width",
    "source_height",
    "pptx_left_emu",
    "pptx_top_emu",
    "pptx_width_emu",
    "pptx_height_emu",
    "text_readability",
    "recognized_text",
    "text_source",
    "text_accuracy_status",
    "converted_to_editable_text",
    "structure_complexity",
    "crop_mask_type",
    "crop_strategy",
    "alpha_strategy",
    "transparent_background",
    "alpha_note",
    "crop_decision",
    "review_status",
    "review_note",
    "needs_manual_review",
    "text_residue",
    "simple_structure_residue",
    "neighboring_content",
    "clipped_content",
    "duplicate_content",
    "overlap_exception",
    "conversion_note",
    "source_page_path",
    "semantic_object_name",
    "semantic_scope",
    "semantic_page_tile_status",
    "semantic_independence_status",
    "semantic_audit_path",
    "background_purity_status",
    "background_scene_complexity",
    "background_origin",
    "foreground_removal_method",
    "background_audit_path",
    "original_full_slide_reuse",
    "icon_residue",
    "shape_residue",
    "chart_residue",
}

QUALITY_REQUIRED = {
    "slide_id",
    "slide_pass",
    "preview_status",
    "text_accuracy_status",
    "aspect_ratio_status",
    "unresolved_review_count",
    "crop_failure_count",
    "duplicate_asset_count",
    "overflow_count",
    "clipping_count",
    "overlap_count",
    "background_gate_status",
    "original_full_slide_status",
    "semantic_split_status",
    "semantic_tile_status",
    "quality_notes",
}

REPRESENTATIONS = {"editable_text", "native_shape", "semantic_png", "background", "backup_only"}
PNG_TRANSPARENCY_TYPES = {"icon", "logo", "symbol", "arrow", "illustration"}
ZERO_COUNTERS = {
    "unresolved_review_count",
    "crop_failure_count",
    "duplicate_asset_count",
    "overflow_count",
    "clipping_count",
    "overlap_count",
}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
SEMANTIC_SCOPES = {"single_object", "single_composite_visual"}
BACKGROUND_PURITY = {"pure_background", "foreground_removed"}
IMAGE_GEN_ORIGINS = {"image_gen_inpainted", "image_gen_reconstructed"}


def read_csv(path: Path) -> List[Dict[str, str]]:
    with path.open("r", newline="", encoding="utf-8-sig") as handle:
        return [dict(row) for row in csv.DictReader(handle)]


def as_float(value: object, default: float = 0.0) -> float:
    try:
        return float(str(value).strip())
    except (TypeError, ValueError):
        return default


def as_int(value: object, default: int = -1) -> int:
    try:
        return int(float(str(value).strip()))
    except (TypeError, ValueError):
        return default


def norm(value: object) -> str:
    return str(value or "").strip().lower()


def is_yes(value: object) -> bool:
    return norm(value) in {"yes", "true", "1", "approved", "pass", "passed"}


def resolve_artifact_path(root: Path, value: object) -> Path | None:
    raw = str(value or "").strip()
    if not raw:
        return None
    path = Path(raw)
    return path if path.is_absolute() else root / path


def coverage_ratio(row: Dict[str, str]) -> float:
    sw = max(1.0, as_float(row.get("source_width"), 1.0))
    sh = max(1.0, as_float(row.get("source_height"), 1.0))
    return max(0.0, as_float(row.get("width"))) * max(0.0, as_float(row.get("height"))) / (sw * sh)


def compare_full_page_images(asset: Path, source: Path) -> Tuple[bool, float, str]:
    """Return whether an asset is an original or near-identical full-page screenshot."""
    try:
        from PIL import Image, ImageChops, ImageStat
    except ImportError:
        return False, 1.0, "Pillow is unavailable"

    try:
        with Image.open(asset) as asset_image, Image.open(source) as source_image:
            candidate = asset_image.convert("RGB")
            reference = source_image.convert("RGB")
            candidate_ratio = candidate.width / max(1, candidate.height)
            reference_ratio = reference.width / max(1, reference.height)
            if abs(candidate_ratio - reference_ratio) / max(reference_ratio, 1e-9) > 0.005:
                return False, 1.0, "ok"
            if candidate.size != reference.size:
                candidate = candidate.resize(reference.size, Image.Resampling.LANCZOS)
            difference = ImageChops.difference(candidate, reference)
            mean_abs = sum(ImageStat.Stat(difference).mean) / (3.0 * 255.0)
            return mean_abs <= 0.003, mean_abs, "ok"
    except Exception as exc:  # pragma: no cover - defensive reporting for corrupt assets
        return False, 1.0, str(exc)


def compare_full_page_blob(asset_blob: bytes, source: Path) -> Tuple[bool, float, str]:
    """Compare an image embedded in the PPTX package with an immutable source page."""
    try:
        from PIL import Image, ImageChops, ImageStat
    except ImportError:
        return False, 1.0, "Pillow is unavailable"

    try:
        with Image.open(io.BytesIO(asset_blob)) as asset_image, Image.open(source) as source_image:
            candidate = asset_image.convert("RGB")
            reference = source_image.convert("RGB")
            candidate_ratio = candidate.width / max(1, candidate.height)
            reference_ratio = reference.width / max(1, reference.height)
            if abs(candidate_ratio - reference_ratio) / max(reference_ratio, 1e-9) > 0.005:
                return False, 1.0, "ok"
            if candidate.size != reference.size:
                candidate = candidate.resize(reference.size, Image.Resampling.LANCZOS)
            difference = ImageChops.difference(candidate, reference)
            mean_abs = sum(ImageStat.Stat(difference).mean) / (3.0 * 255.0)
            return mean_abs <= 0.003, mean_abs, "ok"
    except Exception as exc:  # pragma: no cover - defensive reporting for corrupt assets
        return False, 1.0, str(exc)


def issue(
    issues: List[Dict[str, str]],
    code: str,
    message: str,
    *,
    slide_id: str = "",
    element_id: str = "",
    severity: str = "error",
) -> None:
    issues.append(
        {
            "severity": severity,
            "code": code,
            "slide_id": slide_id,
            "element_id": element_id,
            "message": message,
        }
    )


def require_artifacts(root: Path, pptx: Path, issues: List[Dict[str, str]]) -> None:
    required = [
        root / "source_pages",
        root / "split_png_elements",
        root / "visual_elements_manifest.csv",
        root / "visual_elements_manifest.json",
        root / "image_source_report.csv",
        root / "image_source_report.json",
        root / "review",
        root / "quality_preview",
        root / "quality_report" / "quality_report.csv",
        root / "quality_report" / "quality_report.json",
    ]
    for path in required:
        if not path.exists():
            issue(issues, "missing_artifact", f"Required artifact is missing: {path}")
    if not pptx.exists():
        issue(issues, "missing_pptx", f"Final PPTX is missing: {pptx}")


def validate_manifest(
    root: Path,
    rows: Sequence[Dict[str, str]],
    issues: List[Dict[str, str]],
    overlap_threshold: float,
) -> Tuple[List[Dict[str, str]], List[Dict[str, str]]]:
    if not rows:
        issue(issues, "empty_manifest", "The refined manifest has no rows.")
        return [], []

    missing = sorted(MANIFEST_REQUIRED - set(rows[0]))
    if missing:
        issue(issues, "manifest_fields_missing", "Missing manifest fields: " + ", ".join(missing))

    element_ids: set[str] = set()
    shape_names: set[str] = set()
    owner_ids: set[str] = set()
    semantic_rows: List[Dict[str, str]] = []
    large_background_rows: List[Dict[str, str]] = []

    for row in rows:
        sid = row.get("slide_id", "")
        eid = row.get("element_id", "")
        rep = norm(row.get("final_representation"))
        included = is_yes(row.get("included_in_final"))
        shape_name = str(row.get("pptx_shape_name") or "").strip()

        if not eid or eid in element_ids:
            issue(issues, "element_id_invalid", f"Missing or duplicate element_id: {eid!r}", slide_id=sid, element_id=eid)
        element_ids.add(eid)

        if rep not in REPRESENTATIONS:
            issue(issues, "representation_invalid", f"Invalid final_representation: {rep!r}", slide_id=sid, element_id=eid)

        if included:
            if not shape_name:
                issue(issues, "shape_name_missing", "Included final object has no pptx_shape_name.", slide_id=sid, element_id=eid)
            elif shape_name in shape_names:
                issue(issues, "shape_name_duplicate", f"Duplicate pptx_shape_name: {shape_name}", slide_id=sid, element_id=eid)
            shape_names.add(shape_name)

        if included and norm(row.get("review_status")) not in {"approved", "not_applicable"}:
            issue(issues, "review_unresolved", "Included object is not approved.", slide_id=sid, element_id=eid)
        if included and is_yes(row.get("needs_manual_review")):
            issue(issues, "manual_review_open", "needs_manual_review remains yes.", slide_id=sid, element_id=eid)

        if included and norm(row.get("text_readability")) == "readable" and rep != "editable_text":
            issue(issues, "readable_text_not_editable", "Readable text is not represented as editable_text.", slide_id=sid, element_id=eid)
        if included and norm(row.get("structure_complexity")) == "simple" and rep != "native_shape":
            issue(issues, "simple_structure_not_native", "Simple structure is not represented as native_shape.", slide_id=sid, element_id=eid)

        if rep == "backup_only" and included:
            issue(issues, "backup_visible", "backup_only asset is included in the final deck.", slide_id=sid, element_id=eid)

        if included and rep in {"semantic_png", "background"}:
            rel = str(row.get("relative_path") or "").strip()
            if not rel or not (root / rel).exists():
                issue(issues, "png_missing", f"Final image asset is missing: {rel!r}", slide_id=sid, element_id=eid)
            asset = resolve_artifact_path(root, rel)
            source = resolve_artifact_path(root, row.get("source_page_path"))
            if asset is not None and asset.exists() and source is not None and source.exists():
                near_identical, difference, detail = compare_full_page_images(asset, source)
                if detail != "ok":
                    issue(issues, "full_page_identity_check_failed", f"Could not compare final PNG with source: {detail}", slide_id=sid, element_id=eid)
                elif near_identical:
                    issue(issues, "original_full_slide_reused", f"Visible PNG is original or near-identical to the full source page (normalized difference {difference:.4f}).", slide_id=sid, element_id=eid)

        if included and rep == "semantic_png":
            semantic_rows.append(row)
            for field in (
                "text_residue",
                "simple_structure_residue",
                "neighboring_content",
                "clipped_content",
                "duplicate_content",
                "icon_residue",
                "shape_residue",
                "chart_residue",
            ):
                if norm(row.get(field)) != "no":
                    issue(issues, f"semantic_{field}", f"semantic_png requires {field}=no.", slide_id=sid, element_id=eid)

            if norm(row.get("semantic_scope")) not in SEMANTIC_SCOPES:
                issue(issues, "semantic_scope_invalid", "semantic_png must represent one object or one coherent composite visual.", slide_id=sid, element_id=eid)
            if not str(row.get("semantic_object_name") or "").strip() or norm(row.get("semantic_object_name")) == "not_applicable":
                issue(issues, "semantic_object_unnamed", "semantic_png must name its semantic object.", slide_id=sid, element_id=eid)
            if norm(row.get("semantic_page_tile_status")) != "not_page_tile":
                issue(issues, "semantic_page_tile_unapproved", "semantic_png is not certified as an object-level asset.", slide_id=sid, element_id=eid)
            if norm(row.get("semantic_independence_status")) != "approved":
                issue(issues, "semantic_independence_unapproved", "semantic PNG independence review is not approved.", slide_id=sid, element_id=eid)
            semantic_audit = resolve_artifact_path(root, row.get("semantic_audit_path"))
            if semantic_audit is None or not semantic_audit.exists():
                issue(issues, "semantic_audit_missing", "semantic_png has no existing independence-audit evidence.", slide_id=sid, element_id=eid)

            owner = str(row.get("source_owner_id") or "").strip()
            if not owner or owner in owner_ids:
                issue(issues, "source_owner_invalid", f"Missing or duplicate source_owner_id: {owner!r}", slide_id=sid, element_id=eid)
            owner_ids.add(owner)

            if norm(row.get("crop_decision")) != "keep":
                issue(issues, "crop_not_closed", "Final semantic PNG must have crop_decision=keep.", slide_id=sid, element_id=eid)

            if norm(row.get("element_type")) in PNG_TRANSPARENCY_TYPES and norm(row.get("transparent_background")) != "yes":
                note = norm(row.get("alpha_note"))
                if "approved" not in note:
                    issue(issues, "transparency_not_approved", "Foreground PNG lacks transparency and an approved alpha exception.", slide_id=sid, element_id=eid)

            coverage = coverage_ratio(row)
            if coverage >= 0.85:
                issue(issues, "full_slide_semantic_png", "A semantic PNG covers at least 85% of the slide; split real objects instead of relabeling the page image.", slide_id=sid, element_id=eid)

        if included and rep == "background":
            coverage = coverage_ratio(row)
            if norm(row.get("original_full_slide_reuse")) != "no":
                issue(issues, "original_full_slide_not_disclaimed", "Visible background must set original_full_slide_reuse=no.", slide_id=sid, element_id=eid)

            source = resolve_artifact_path(root, row.get("source_page_path"))
            if coverage >= 0.85:
                if source is None or not source.exists():
                    issue(issues, "source_page_missing_for_identity_check", "Large background has no existing immutable source_page_path.", slide_id=sid, element_id=eid)

            if coverage > 0.95:
                large_background_rows.append(row)
                purity = norm(row.get("background_purity_status"))
                complexity = norm(row.get("background_scene_complexity"))
                origin = norm(row.get("background_origin"))
                method = norm(row.get("foreground_removal_method"))

                if purity not in BACKGROUND_PURITY:
                    issue(issues, "background_purity_unproven", "Background over 95% coverage must be pure or foreground-removed.", slide_id=sid, element_id=eid)
                if complexity == "simple":
                    issue(issues, "simple_background_not_native", "Simple full-page backgrounds must be rebuilt with native PowerPoint shapes.", slide_id=sid, element_id=eid)
                if complexity not in {"complex", "simple"}:
                    issue(issues, "background_complexity_missing", "Large background must declare simple or complex scene complexity.", slide_id=sid, element_id=eid)
                if purity == "foreground_removed" and (method != "image_gen" or origin not in IMAGE_GEN_ORIGINS):
                    issue(issues, "complex_background_not_image_gen", "Foreground removal on a complex large background requires Image Gen and an Image Gen origin.", slide_id=sid, element_id=eid)
                if origin not in IMAGE_GEN_ORIGINS | {"derived_background"}:
                    issue(issues, "background_origin_invalid", "Large background has no approved non-screenshot origin.", slide_id=sid, element_id=eid)
                for field in ("text_residue", "icon_residue", "shape_residue", "chart_residue"):
                    if norm(row.get(field)) != "no":
                        issue(issues, f"background_{field}", f"Large background requires {field}=no.", slide_id=sid, element_id=eid)
                audit = resolve_artifact_path(root, row.get("background_audit_path"))
                if audit is None or not audit.exists():
                    issue(issues, "background_audit_missing", "Background over 95% coverage has no existing source-cleaned-difference evidence.", slide_id=sid, element_id=eid)

    validate_semantic_overlap(semantic_rows, issues, overlap_threshold)
    validate_semantic_collections(semantic_rows, issues)
    return semantic_rows, large_background_rows


def semantic_union_ratio(rows: Sequence[Dict[str, str]]) -> float:
    if not rows:
        return 0.0
    source_width = max(1.0, as_float(rows[0].get("source_width"), 1.0))
    source_height = max(1.0, as_float(rows[0].get("source_height"), 1.0))
    xs = sorted({as_float(row.get("x")) for row in rows} | {as_float(row.get("x")) + max(0.0, as_float(row.get("width"))) for row in rows})
    area = 0.0
    for left, right in zip(xs, xs[1:]):
        if right <= left:
            continue
        intervals = []
        for row in rows:
            x = as_float(row.get("x"))
            width = max(0.0, as_float(row.get("width")))
            if x < right and x + width > left:
                y = as_float(row.get("y"))
                intervals.append((y, y + max(0.0, as_float(row.get("height")))))
        intervals.sort()
        covered_y = 0.0
        if intervals:
            start, end = intervals[0]
            for next_start, next_end in intervals[1:]:
                if next_start <= end:
                    end = max(end, next_end)
                else:
                    covered_y += max(0.0, end - start)
                    start, end = next_start, next_end
            covered_y += max(0.0, end - start)
        area += (right - left) * covered_y
    return area / (source_width * source_height)


def validate_semantic_collections(rows: Sequence[Dict[str, str]], issues: List[Dict[str, str]]) -> None:
    by_slide: Dict[str, List[Dict[str, str]]] = defaultdict(list)
    for row in rows:
        by_slide[str(row.get("slide_id") or "")].append(row)

    for sid, slide_rows in by_slide.items():
        if len(slide_rows) == 1:
            issue(issues, "semantic_split_insufficient", "A slide using semantic PNGs must contain multiple independent semantic assets, not one residual screenshot.", slide_id=sid, element_id=str(slide_rows[0].get("element_id") or ""))
        rectangular = [row for row in slide_rows if norm(row.get("crop_mask_type")) in {"rectangle", "full_background"}]
        union_ratio = semantic_union_ratio(slide_rows)
        if len(slide_rows) >= 2 and len(rectangular) == len(slide_rows) and union_ratio >= 0.85:
            issue(issues, "semantic_page_tiling_detected", f"Rectangular semantic crops cover {union_ratio:.1%} of the page and form a page-tile reconstruction.", slide_id=sid)


def validate_semantic_overlap(
    rows: Sequence[Dict[str, str]],
    issues: List[Dict[str, str]],
    threshold: float,
) -> None:
    for index, left in enumerate(rows):
        for right in rows[index + 1 :]:
            if left.get("slide_id") != right.get("slide_id"):
                continue
            lx, ly = as_float(left.get("x")), as_float(left.get("y"))
            lw, lh = as_float(left.get("width")), as_float(left.get("height"))
            rx, ry = as_float(right.get("x")), as_float(right.get("y"))
            rw, rh = as_float(right.get("width")), as_float(right.get("height"))
            iw = max(0.0, min(lx + lw, rx + rw) - max(lx, rx))
            ih = max(0.0, min(ly + lh, ry + rh) - max(ly, ry))
            smaller = min(max(1.0, lw * lh), max(1.0, rw * rh))
            ratio = iw * ih / smaller
            if ratio >= threshold and not (is_yes(left.get("overlap_exception")) or is_yes(right.get("overlap_exception"))):
                issue(
                    issues,
                    "semantic_overlap_unapproved",
                    f"Semantic ownership boxes overlap by {ratio:.1%} without an approved exception: {right.get('element_id')}",
                    slide_id=str(left.get("slide_id") or ""),
                    element_id=str(left.get("element_id") or ""),
                )


def validate_crop_audit(
    root: Path,
    semantic_rows: Sequence[Dict[str, str]],
    issues: List[Dict[str, str]],
) -> None:
    if not semantic_rows:
        return
    path = root / "review" / "crop_visual_audit" / "crop_audit.csv"
    if not path.exists():
        issue(issues, "crop_audit_missing", f"Crop audit is required for semantic PNG assets: {path}")
        return
    rows = read_csv(path)
    required = {"element_id", "crop_decision", "review_status"}
    if rows and not required.issubset(rows[0]):
        issue(issues, "crop_audit_fields_missing", "crop_audit.csv must include element_id, crop_decision, and review_status.")
    by_id = {str(row.get("element_id") or ""): row for row in rows}
    for manifest_row in semantic_rows:
        eid = str(manifest_row.get("element_id") or "")
        audit = by_id.get(eid)
        if audit is None:
            issue(issues, "crop_audit_row_missing", "Semantic PNG has no crop audit row.", slide_id=str(manifest_row.get("slide_id") or ""), element_id=eid)
            continue
        if norm(audit.get("crop_decision")) != "keep" or norm(audit.get("review_status")) != "approved":
            issue(issues, "crop_audit_unresolved", "Crop audit row is not closed and approved.", slide_id=str(manifest_row.get("slide_id") or ""), element_id=eid)


def validate_background_audit(
    root: Path,
    background_rows: Sequence[Dict[str, str]],
    issues: List[Dict[str, str]],
) -> None:
    if not background_rows:
        return
    path = root / "review" / "background_visual_audit" / "background_audit.csv"
    if not path.exists():
        issue(issues, "background_audit_csv_missing", f"Large-background audit is required: {path}")
        return
    rows = read_csv(path)
    required = {
        "element_id",
        "coverage_ratio",
        "background_purity_status",
        "background_origin",
        "foreground_removal_method",
        "original_full_slide_reuse",
        "text_residue",
        "icon_residue",
        "shape_residue",
        "chart_residue",
        "review_status",
        "evidence_file",
    }
    if not rows:
        issue(issues, "background_audit_empty", "background_audit.csv has no rows.")
        return
    missing = sorted(required - set(rows[0]))
    if missing:
        issue(issues, "background_audit_fields_missing", "Missing background audit fields: " + ", ".join(missing))
    by_id = {str(row.get("element_id") or ""): row for row in rows}
    for manifest_row in background_rows:
        sid = str(manifest_row.get("slide_id") or "")
        eid = str(manifest_row.get("element_id") or "")
        audit = by_id.get(eid)
        if audit is None:
            issue(issues, "background_audit_row_missing", "Large background has no audit row.", slide_id=sid, element_id=eid)
            continue
        if as_float(audit.get("coverage_ratio")) <= 0.95:
            issue(issues, "background_audit_coverage_invalid", "Large-background audit coverage must be greater than 0.95.", slide_id=sid, element_id=eid)
        if norm(audit.get("background_purity_status")) not in BACKGROUND_PURITY:
            issue(issues, "background_audit_purity_failed", "Background audit does not prove purity or foreground removal.", slide_id=sid, element_id=eid)
        if norm(audit.get("background_purity_status")) == "foreground_removed":
            if norm(audit.get("foreground_removal_method")) != "image_gen" or norm(audit.get("background_origin")) not in IMAGE_GEN_ORIGINS:
                issue(issues, "background_audit_image_gen_failed", "Audit does not prove Image Gen foreground removal for the complex background.", slide_id=sid, element_id=eid)
        if norm(audit.get("original_full_slide_reuse")) != "no":
            issue(issues, "background_audit_original_reuse", "Background audit does not rule out original full-slide reuse.", slide_id=sid, element_id=eid)
        for field in ("text_residue", "icon_residue", "shape_residue", "chart_residue"):
            if norm(audit.get(field)) != "no":
                issue(issues, f"background_audit_{field}", f"Background audit requires {field}=no.", slide_id=sid, element_id=eid)
        if norm(audit.get("review_status")) != "approved":
            issue(issues, "background_audit_unapproved", "Large-background audit is not approved.", slide_id=sid, element_id=eid)
        evidence = resolve_artifact_path(root, audit.get("evidence_file"))
        if evidence is None or not evidence.exists():
            issue(issues, "background_audit_evidence_missing", "Background audit evidence file is missing.", slide_id=sid, element_id=eid)


def validate_quality_report(root: Path, issues: List[Dict[str, str]]) -> None:
    path = root / "quality_report" / "quality_report.csv"
    if not path.exists():
        return
    rows = read_csv(path)
    if not rows:
        issue(issues, "quality_report_empty", "quality_report.csv has no slide rows.")
        return
    missing = sorted(QUALITY_REQUIRED - set(rows[0]))
    if missing:
        issue(issues, "quality_fields_missing", "Missing quality report fields: " + ", ".join(missing))
    for row in rows:
        sid = row.get("slide_id", "")
        if not is_yes(row.get("slide_pass")):
            issue(issues, "slide_not_passed", "slide_pass is not yes.", slide_id=sid)
        if norm(row.get("preview_status")) != "approved":
            issue(issues, "preview_not_approved", "preview_status must be approved.", slide_id=sid)
        if norm(row.get("text_accuracy_status")) not in {"exact", "approved"}:
            issue(issues, "text_accuracy_failed", "text_accuracy_status must be exact or approved.", slide_id=sid)
        if norm(row.get("aspect_ratio_status")) != "match":
            issue(issues, "aspect_report_failed", "aspect_ratio_status must be match.", slide_id=sid)
        if norm(row.get("background_gate_status")) not in {"approved", "not_applicable"}:
            issue(issues, "background_gate_report_failed", "background_gate_status must be approved or not_applicable.", slide_id=sid)
        if norm(row.get("original_full_slide_status")) != "not_reused":
            issue(issues, "original_full_slide_report_failed", "original_full_slide_status must be not_reused.", slide_id=sid)
        if norm(row.get("semantic_split_status")) not in {"approved", "not_applicable"}:
            issue(issues, "semantic_split_report_failed", "semantic_split_status must be approved or not_applicable.", slide_id=sid)
        if norm(row.get("semantic_tile_status")) not in {"not_page_tiles", "not_applicable"}:
            issue(issues, "semantic_tile_report_failed", "semantic_tile_status must be not_page_tiles or not_applicable.", slide_id=sid)
        for field in ZERO_COUNTERS:
            if as_int(row.get(field)) != 0:
                issue(issues, "quality_counter_nonzero", f"{field} must be 0, got {row.get(field)!r}.", slide_id=sid)


def ordered_slide_ids(rows: Iterable[Dict[str, str]]) -> List[str]:
    result: List[str] = []
    seen: set[str] = set()
    for row in rows:
        sid = str(row.get("slide_id") or "")
        if sid and sid not in seen:
            result.append(sid)
            seen.add(sid)
    return result


def validate_pptx(
    pptx_path: Path,
    root: Path,
    rows: Sequence[Dict[str, str]],
    issues: List[Dict[str, str]],
    aspect_tolerance: float,
) -> None:
    if not pptx_path.exists():
        return
    try:
        from pptx import Presentation
    except ImportError:
        issue(issues, "python_pptx_missing", "python-pptx is required to validate the final deck.")
        return

    prs = Presentation(pptx_path)
    slide_ids = ordered_slide_ids(rows)
    if len(prs.slides) != len(slide_ids):
        issue(issues, "slide_count_mismatch", f"PPTX has {len(prs.slides)} slides but manifest has {len(slide_ids)} slide IDs.")

    expected_names = {
        str(row.get("pptx_shape_name") or "").strip()
        for row in rows
        if is_yes(row.get("included_in_final")) and str(row.get("pptx_shape_name") or "").strip()
    }
    actual_names: set[str] = set()
    overflow = 0
    tolerance_x = int(prs.slide_width * 0.001)
    tolerance_y = int(prs.slide_height * 0.001)
    for slide in prs.slides:
        for shape in slide.shapes:
            actual_names.add(shape.name)
            if (
                shape.left < -tolerance_x
                or shape.top < -tolerance_y
                or shape.left + shape.width > prs.slide_width + tolerance_x
                or shape.top + shape.height > prs.slide_height + tolerance_y
            ):
                overflow += 1
    if overflow:
        issue(issues, "pptx_overflow", f"Found {overflow} out-of-bounds PowerPoint objects.")

    missing_names = sorted(expected_names - actual_names)
    extra_names = sorted(actual_names - expected_names)
    if missing_names:
        issue(issues, "pptx_shapes_missing", "Manifest shape names missing from PPTX: " + ", ".join(missing_names[:20]))
    if extra_names:
        issue(issues, "pptx_shapes_unmanifested", "PPTX shapes missing from manifest: " + ", ".join(extra_names[:20]))

    by_slide: Dict[str, Dict[str, str]] = {}
    for row in rows:
        by_slide.setdefault(str(row.get("slide_id") or ""), row)
    pptx_ratio = prs.slide_width / max(1, prs.slide_height)
    for sid in slide_ids:
        row = by_slide[sid]
        source_ratio = as_float(row.get("source_width")) / max(1.0, as_float(row.get("source_height"), 1.0))
        relative_error = abs(pptx_ratio - source_ratio) / max(source_ratio, 1e-9)
        if relative_error > aspect_tolerance:
            issue(issues, "aspect_ratio_mismatch", f"PPTX/source aspect ratio error is {relative_error:.3%}.", slide_id=sid)

    validate_pptx_embedded_images(pptx_path, root, rows, issues)


def validate_pptx_embedded_images(
    pptx_path: Path,
    root: Path,
    rows: Sequence[Dict[str, str]],
    issues: List[Dict[str, str]],
) -> None:
    source_paths: List[Path] = []
    seen: set[Path] = set()
    for row in rows:
        source = resolve_artifact_path(root, row.get("source_page_path"))
        if source is not None and source.exists():
            source = source.resolve()
            if source not in seen:
                seen.add(source)
                source_paths.append(source)
    if not source_paths:
        return

    try:
        with zipfile.ZipFile(pptx_path) as package:
            media_names = [
                name
                for name in package.namelist()
                if name.startswith("ppt/media/") and Path(name).suffix.lower() in IMAGE_SUFFIXES
            ]
            for media_name in media_names:
                blob = package.read(media_name)
                for source in source_paths:
                    near_identical, difference, detail = compare_full_page_blob(blob, source)
                    if detail != "ok":
                        issue(issues, "pptx_media_identity_check_failed", f"Could not inspect {media_name}: {detail}")
                        break
                    if near_identical:
                        issue(issues, "pptx_original_full_slide_media", f"PPTX embeds an original or near-identical full source screenshot in {media_name} (normalized difference {difference:.4f}).")
                        break
    except (OSError, zipfile.BadZipFile) as exc:
        issue(issues, "pptx_media_scan_failed", f"Could not scan embedded PPTX media: {exc}")


def write_report(root: Path, issues: Sequence[Dict[str, str]], pptx: Path) -> None:
    report_dir = root / "quality_report"
    report_dir.mkdir(parents=True, exist_ok=True)
    errors = [item for item in issues if item["severity"] == "error"]
    payload = {
        "status": "passed" if not errors else "failed",
        "error_count": len(errors),
        "warning_count": len(issues) - len(errors),
        "output_root": str(root),
        "pptx": str(pptx),
        "issues": list(issues),
    }
    (report_dir / "delivery_validation.json").write_text(
        json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    with (report_dir / "delivery_validation.csv").open("w", newline="", encoding="utf-8-sig") as handle:
        writer = csv.DictWriter(handle, fieldnames=["severity", "code", "slide_id", "element_id", "message"])
        writer.writeheader()
        writer.writerows(issues)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Validate a v2.1 refined editable PowerPoint delivery.")
    parser.add_argument("output_dir", type=Path, help="Refined editable output directory")
    parser.add_argument("--pptx", type=Path, required=True, help="Final refined PPTX")
    parser.add_argument("--aspect-tolerance", type=float, default=0.001, help="Maximum relative aspect-ratio error")
    parser.add_argument("--overlap-threshold", type=float, default=0.10, help="Intersection area divided by smaller semantic bbox")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    root = args.output_dir.expanduser().resolve()
    pptx = args.pptx.expanduser().resolve()
    issues: List[Dict[str, str]] = []
    require_artifacts(root, pptx, issues)

    manifest_path = root / "visual_elements_manifest.csv"
    rows = read_csv(manifest_path) if manifest_path.exists() else []
    semantic_rows, large_background_rows = validate_manifest(root, rows, issues, args.overlap_threshold)
    validate_crop_audit(root, semantic_rows, issues)
    validate_background_audit(root, large_background_rows, issues)
    validate_quality_report(root, issues)
    validate_pptx(pptx, root, rows, issues, args.aspect_tolerance)
    write_report(root, issues, pptx)

    errors = [item for item in issues if item["severity"] == "error"]
    if errors:
        print(f"Delivery validation failed with {len(errors)} error(s).")
        print(root / "quality_report" / "delivery_validation.csv")
        return 1
    print("Delivery validation passed.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
